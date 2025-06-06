import os
import json
import re
from django.conf import settings
from openai import OpenAI
from docx import Document as DocxDocument
from pptx import Presentation as PptxPresentation
from PyPDF2 import PdfReader

client = OpenAI(
    base_url="https://openrouter.ai/api/v1",
    api_key=settings.OPENROUTER_API_KEY,
)

# ===== UTILITY FUNCTIONS =====

def extract_json_from_response(text):
    """Extract JSON from markdown code blocks or return as-is"""
    # Try to find JSON in code blocks first
    
    # Look for ```json ... ``` blocks
    json_match = re.search(r'```json\s*\n(.*?)\n```', text, re.DOTALL)
    if json_match:
        print("🔍 Found JSON in ```json code block")
        return json_match.group(1).strip()
    
    # Look for ``` ... ``` blocks (any language)
    code_match = re.search(r'```[a-zA-Z]*\s*\n(.*?)\n```', text, re.DOTALL)
    if code_match:
        print("🔍 Found JSON in generic code block")
        return code_match.group(1).strip()
    
    # Look for just the JSON object (starts with { and ends with })
    json_obj_match = re.search(r'\{.*\}', text, re.DOTALL)
    if json_obj_match:
        print("🔍 Found JSON object in text")
        return json_obj_match.group(0).strip()
    
    # Return as-is if no patterns found
    print("🔍 No JSON patterns found, returning raw text")
    return text.strip()


def validate_and_improve_title(title, content_type, material_title):
    """Validate and improve AI-generated titles"""
    title = title.strip()
    
    # List of generic titles to avoid
    generic_titles = {
        'flashcards': ['flashcards', 'cards', 'study cards', 'review cards', 'flash cards'],
        'quiz': ['quiz', 'test', 'exam', 'questions', 'practice quiz'],
        'notes': ['notes', 'study notes', 'study guide', 'summary', 'note']
    }
    
    if not title or title.lower() in generic_titles.get(content_type, []):
        # Generate a better title
        from datetime import datetime
        timestamp = datetime.now().strftime("%B %d")
        return f"{material_title} - {content_type.title()} ({timestamp})"
    
    # Ensure title is descriptive enough
    if len(title.split()) < 2:
        return f"{title} - {material_title}"
    
    return title


def generate_fallback_title(content_type, material_title, timestamp=None):
    """Generate a fallback title if AI doesn't provide a good one"""
    from datetime import datetime
    
    if timestamp is None:
        timestamp = datetime.now().strftime("%m/%d")
    
    fallbacks = {
        'flashcards': f"{material_title} - Flashcards ({timestamp})",
        'quiz': f"{material_title} - Quiz ({timestamp})", 
        'notes': f"{material_title} - Study Notes ({timestamp})"
    }
    
    return fallbacks.get(content_type, f"{material_title} - {content_type.title()} ({timestamp})")

# ===== FILE EXTRACTION FUNCTIONS =====

def extract_text_from_pdf(path_on_disk: str) -> str:
    reader = PdfReader(path_on_disk)
    full_text = []
    for page in reader.pages:
        full_text.append(page.extract_text() or "")
    return "\n".join(full_text)

def extract_text_from_docx(path_on_disk: str) -> str:
    doc = DocxDocument(path_on_disk)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paragraphs)

def read_text_file(path_on_disk: str) -> str:
    with open(path_on_disk, "r", encoding="utf-8") as f:
        return f.read()

def extract_text_from_pptx(path_on_disk: str) -> str:
    prs = PptxPresentation(path_on_disk)
    slides_text = []
    for slide in prs.slides:
        slide_paras = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_paras.append(shape.text.strip())
        if slide_paras:
            slides_text.append("\n".join(slide_paras))
    return "\n\n".join(slides_text)

def gather_material_text(material, specific_attachment_ids=None) -> str:
    """Extract text from specific attachments or all attachments in a material"""
    
    # ✅ Force refresh from database to get latest attachments
    material.refresh_from_db()
    
    if specific_attachment_ids:
        attachments = material.attachments.filter(id__in=specific_attachment_ids)
        print(f"🎯 Processing {len(specific_attachment_ids)} specific attachments from material '{material.title}'")
    else:
        attachments = material.attachments.all()
        print(f"📁 Processing all {attachments.count()} attachments from material '{material.title}'")
    
    attachment_count = attachments.count()
    
    # ✅ Better error message if no attachments
    if attachment_count == 0:
        if specific_attachment_ids:
            raise ValueError(
                f"No attachments found with IDs {specific_attachment_ids} for material '{material.title}'. "
                "The specified files may not exist or may still be uploading."
            )
        else:
            raise ValueError(
                f"No attachments found for material '{material.title}'. "
                "Files may still be uploading. Please wait a moment and try again."
            )
    
    texts = []
    supported_extensions = ['.pdf', '.docx', '.txt', '.pptx']
    
    for attachment in attachments:
        try:
            path = attachment.file.path
            ext = os.path.splitext(path)[1].lower()
            
            print(f"🔍 Processing: {attachment.file.name} (ext: {ext})")
            
            if ext not in supported_extensions:
                print(f"⚠️ Skipping unsupported file type: {ext}")
                continue
                
            if not os.path.exists(path):
                print(f"❌ File not found on disk: {path}")
                continue
            
            extracted_text = ""
            if ext == ".pdf":
                extracted_text = extract_text_from_pdf(path)
            elif ext == ".docx":
                extracted_text = extract_text_from_docx(path)
            elif ext == ".txt":
                extracted_text = read_text_file(path)
            elif ext == ".pptx":
                extracted_text = extract_text_from_pptx(path)
            
            if extracted_text and extracted_text.strip():
                texts.append(extracted_text)
                print(f"✅ Extracted {len(extracted_text)} characters from {attachment.file.name}")
            else:
                print(f"⚠️ No text content in {attachment.file.name}")
                
        except Exception as e:
            print(f"❌ Error processing {attachment.file.name}: {str(e)}")
            continue
    
    result = "\n\n".join(texts).strip()
    
    # ✅ Better validation of extracted text
    if not result:
        if specific_attachment_ids:
            file_names = [att.file.name for att in attachments]
            raise ValueError(
                f"No readable text found in specified files: {', '.join(file_names)}. "
                "Please ensure files contain text content and are in supported formats (PDF, DOCX, TXT, PPTX)."
            )
        else:
            file_names = [att.file.name for att in attachments]
            raise ValueError(
                f"No readable text found in uploaded files: {', '.join(file_names)}. "
                "Please ensure files contain text content and are in supported formats (PDF, DOCX, TXT, PPTX)."
            )
    
    print(f"✅ Total extracted text: {len(result)} characters")
    return result

def get_relevant_material_chunks(material, user_prompt, max_chunks=3):
    """
    Instead of sending entire material text, extract only relevant chunks
    based on user's current question for large documents.
    """
    full_text = gather_material_text(material)
    if len(full_text) < 2000:  # Small documents: use full text
        return full_text
    
    # Split into chunks
    chunks = [full_text[i:i+1000] for i in range(0, len(full_text), 1000)]
    
    # For now, just return first few chunks
    # TODO: Could make this smarter with AI-based relevance scoring
    return "\n\n".join(chunks[:max_chunks])

# ===== CONVERSATION SUMMARY FUNCTIONS =====

def generate_conversation_summary(messages, existing_summary=""):
    """
    Generate a concise summary of conversation messages.
    This helps maintain context while reducing token usage.
    """
    if not messages:
        return existing_summary

    # Format messages for summarization
    formatted_messages = []
    for msg in messages:
        role = msg.get('role', 'user')
        content = msg.get('content', '')
        formatted_messages.append(f"{role}: {content}")
    
    messages_text = "\n".join(formatted_messages)
    
    system_prompt = (
        "You are an AI assistant that creates concise conversation summaries. "
        "Summarize the key topics, questions, and important information discussed. "
        "Keep it under 200 words and focus on context that would be helpful for continuing the conversation.\n\n"
    )
    
    if existing_summary:
        system_prompt += f"Previous summary: {existing_summary}\n\n"
        system_prompt += "Update this summary with the new messages below:\n"
    else:
        system_prompt += "Create a summary of this conversation:\n"
    
    system_prompt += f"\nMessages:\n{messages_text}"

    try:
        response = client.chat.completions.create(
            model="deepseek/deepseek-chat-v3-0324:free",
            messages=[{"role": "system", "content": system_prompt}],
            max_tokens=300  # Limit summary length
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        print(f"Failed to generate conversation summary: {e}")
        return existing_summary

def generate_enhanced_conversation_summary(conversation):
    """
    Generate summary with additional metadata for better context
    """
    messages = conversation.messages
    if not messages:
        return {"summary": "", "main_topic": "general", "key_concepts": []}
    
    # Generate summary with structured output
    system_prompt = f"""
    Analyze this conversation and return a JSON summary:
    {{
        "summary": "Brief summary of the conversation (max 150 words)",
        "main_topic": "Primary topic discussed",
        "key_concepts": ["concept1", "concept2", "concept3"],
        "learning_goals": ["goal1", "goal2"],
        "last_context": "What the user was asking about most recently"
    }}
    
    Conversation:
    {conversation.get_context_for_ai()}
    """
    
    try:
        response = client.chat.completions.create(
            model="deepseek/deepseek-chat-v3-0324:free",
            messages=[{"role": "system", "content": system_prompt}],
            max_tokens=400
        )
        
        return json.loads(response.choices[0].message.content)
    except:
        # Fallback to simple summary
        topic = conversation.detect_conversation_topic()
        simple_summary = generate_conversation_summary(messages)
        return {
            "summary": simple_summary, 
            "main_topic": topic,
            "key_concepts": [],
            "learning_goals": [],
            "last_context": ""
        }

# ===== SMART AI RESPONSE FUNCTIONS =====

def generate_ai_response_with_context(conversation, prompt):
    """
    Generate AI response using efficient context management.
    Uses summarized context + recent messages instead of full history.
    """
    material = conversation.material
    
    # Build the system prompt based on conversation context
    system_prompt = conversation.get_context_aware_system_prompt()
    
    # Build the user prompt with context
    prompt_parts = []
    
    # Add material text if it's relevant to the question
    if material and conversation.should_include_material_context(prompt):
        if material.attachments.exists():
            material_text = get_relevant_material_chunks(material, prompt)
            if material_text:
                prompt_parts.append(f"Study Material:\n{material_text}")
    
    # Add conversation context (summary + recent messages)
    conversation_context = conversation.get_context_for_ai()
    if conversation_context:
        prompt_parts.append(f"Conversation Context:\n{conversation_context}")
    
    # Add current user prompt
    prompt_parts.append(f"Current Question:\n{prompt}")
    
    combined_prompt = "\n\n".join(prompt_parts)
    
    try:
        response = client.chat.completions.create(
            model="deepseek/deepseek-chat-v3-0324:free",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": combined_prompt}
            ],
        )
        return response.choices[0].message.content
    except Exception as e:
        raise ValueError(f"AI response generation failed: {str(e)}")

# ===== LEGACY FUNCTIONS (for backward compatibility) =====

def generate_ai_response(text: str) -> str:
    """Simple AI response for basic prompts without conversation context"""
    try:
        response = client.chat.completions.create(
            model="deepseek/deepseek-chat-v3-0324:free",
            messages=[{"role": "user", "content": text}],
        )
        return response.choices[0].message.content
    except Exception as e:
        raise ValueError(f"AI response generation failed: {str(e)}")

def generate_ai_response_for_material(material, prompt: str) -> str:
    """
    Legacy function - kept for backward compatibility.
    For conversations, use generate_ai_response_with_context instead.
    """
    text_body = gather_material_text(material)
    if not text_body:
        raise ValueError("No extractable text found in this Material's attachments.")

    combined = f"{text_body}\n\nUser prompt:\n{prompt}"

    try:
        response = client.chat.completions.create(
            model="deepseek/deepseek-chat-v3-0324:free",
            messages=[{"role": "user", "content": combined}]
        )
        return response.choices[0].message.content
    except Exception as e:
        raise ValueError(f"AI response generation failed: {str(e)}")

# ===== CONTENT GENERATION FUNCTIONS =====

def generate_flashcards_from_material(material, num_cards: int = 5, specific_attachment_ids=None) -> dict:
    """Generate flashcards with more specific titles"""
    text_body = gather_material_text(material, specific_attachment_ids)
    if not text_body:
        if specific_attachment_ids:
            raise ValueError("No extractable text found in the specified attachments.")
        else:
            raise ValueError("No extractable text found in this Material's attachments.")

    # ✅ IMPROVED: More specific prompt for unique titles
    system_prompt = (
        f"You are an AI study helper. From the material below, generate exactly {num_cards} "
        "simple question-answer flashcards. Create a SPECIFIC title that includes key topics "
        "or concepts from the material (not just generic titles like 'Flashcards' or 'Study Cards').\n\n"
        
        "**TITLE REQUIREMENTS:**\n"
        "- Include specific subject matter or main concepts\n"
        "- Add context like chapter/topic if apparent\n"
        "- Use descriptive words, not just 'Flashcards'\n"
        "- Examples: 'Cell Biology: Mitosis & Meiosis', 'Chapter 3: Market Economics', 'Python Functions & Loops'\n\n"
        
        "**RETURN ONLY RAW JSON - NO MARKDOWN, NO CODE BLOCKS, NO EXTRA TEXT.**\n"
        "Start your response immediately with { and end with }\n\n"
        "Expected JSON structure:\n"
        "{\n"
        "  \"title\": \"<specific descriptive title with main concepts>\",\n"
        "  \"description\": \"<description mentioning key topics covered>\",\n"
        "  \"flashcards\": [\n"
        "    { \"question\": \"...\", \"answer\": \"...\" },\n"
        f"    ... (exactly {num_cards} flashcards)\n"
        "  ]\n"
        "}\n\n"
        f"Material:\n\"\"\"\n{text_body}\n\"\"\"\n"
    )

    try:
        response = client.chat.completions.create(
            model="deepseek/deepseek-chat-v3-0324:free",
            messages=[{"role": "system", "content": system_prompt}],
        )
        raw = response.choices[0].message.content.strip()

        # ✅ Extract JSON from response (handles markdown code blocks)
        json_content = extract_json_from_response(raw)
        print(f"🔍 Extracted JSON content: {json_content[:200]}...")

        try:
            parsed = json.loads(json_content)
        except json.JSONDecodeError as e:
            print(f"❌ JSON parsing failed. Raw response: {raw}")
            raise ValueError(f"AI did not return valid JSON for flashcards: {e}\n\nRaw response: {raw}")

        title = parsed.get("title")
        description = parsed.get("description")
        cards = parsed.get("flashcards")

        if not isinstance(title, str) or not title.strip():
            raise ValueError("AI did not return a valid 'title' for the flashcard set.")
        if not isinstance(description, str):
            description = ""
        if not isinstance(cards, list) or len(cards) < 1:
            raise ValueError("AI JSON did not contain a valid 'flashcards' array.")

        validated = []
        for idx, obj in enumerate(cards):
            q = obj.get("question", "").strip()
            a = obj.get("answer", "").strip()
            if not q or not a:
                raise ValueError(f"Flashcard #{idx+1} missing question or answer: {obj}")
            validated.append({"question": q, "answer": a})

        # ✅ Validate and improve title
        improved_title = validate_and_improve_title(title.strip(), 'flashcards', material.title)

        return {
            "title": improved_title,
            "description": description.strip(),
            "flashcards": validated,
        }
    except Exception as e:
        raise ValueError(f"Flashcard generation failed: {str(e)}")


def generate_notes_from_material(material, specific_attachment_ids=None) -> dict:
    """Generate notes with more specific titles"""
    text_body = gather_material_text(material, specific_attachment_ids)
    if not text_body:
        if specific_attachment_ids:
            raise ValueError("No extractable text found in the specified attachments.")
        else:
            raise ValueError("No extractable text found in this Material's attachments.")

    # ✅ IMPROVED: More specific prompt for unique titles
    system_prompt = (
        "You are an AI study helper. Read the material below and produce a comprehensive, well-structured study note "
        "that captures all the core concepts, key points, and important details. Create a SPECIFIC title that "
        "reflects the actual content covered (not generic titles like 'Notes' or 'Study Guide').\n\n"
        
        "**TITLE REQUIREMENTS:**\n"
        "- Include main topics/concepts covered in the material\n"
        "- Add context like subject area, chapter, or theme\n"
        "- Use descriptive words that identify the content\n"
        "- Examples: 'Renaissance Art: Key Artists & Techniques', 'Organic Chemistry: Functional Groups Overview', 'Chapter 8: Supply & Demand Analysis'\n\n"
        
        "**RETURN ONLY RAW JSON - NO MARKDOWN, NO CODE BLOCKS, NO EXTRA TEXT.**\n"
        "Start your response immediately with { and end with }\n\n"
        "Expected JSON structure:\n"
        "{\n"
        "  \"title\": \"<specific descriptive title reflecting main concepts>\",\n"
        "  \"description\": \"<brief description of topics and scope covered>\",\n"
        "  \"content\": \"<comprehensive note content with proper formatting>\"\n"
        "}\n\n"
        "For the content field, create a well-organized note that includes:\n"
        "- Main topics and subtopics with clear headers\n"
        "- Key concepts and definitions\n"
        "- Important details and explanations\n"
        "- Examples or case studies if present\n"
        "- Use markdown formatting (## headers, bullet points, **bold**, etc.)\n"
        "- Make it comprehensive and study-friendly\n\n"
        f"Material:\n\"\"\"\n{text_body}\n\"\"\"\n"
    )

    try:
        response = client.chat.completions.create(
            model="deepseek/deepseek-chat-v3-0324:free",
            messages=[{"role": "system", "content": system_prompt}],
        )
        raw = response.choices[0].message.content.strip()

        # ✅ Extract JSON from response (handles markdown code blocks)
        json_content = extract_json_from_response(raw)
        print(f"🔍 Extracted JSON content: {json_content[:200]}...")

        try:
            parsed = json.loads(json_content)
        except json.JSONDecodeError as e:
            print(f"❌ JSON parsing failed. Raw response: {raw}")
            raise ValueError(f"AI did not return valid JSON for notes: {e}\n\nRaw response: {raw}")

        title = parsed.get("title")
        description = parsed.get("description")
        content = parsed.get("content")

        if not isinstance(title, str) or not title.strip():
            raise ValueError("AI did not return a valid 'title' for the note.")
        if not isinstance(description, str):
            description = ""
        if not isinstance(content, str) or not content.strip():
            raise ValueError("AI did not return valid 'content' for the note.")

        # ✅ Validate and improve title
        improved_title = validate_and_improve_title(title.strip(), 'notes', material.title)

        return {
            "title": improved_title,
            "description": description.strip(),
            "content": content.strip(),
        }
    except Exception as e:
        raise ValueError(f"Notes generation failed: {str(e)}")


def generate_quiz_from_material(material, num_questions: int = 5, specific_attachment_ids=None) -> dict:
    """Generate quiz with more specific titles"""
    text_body = gather_material_text(material, specific_attachment_ids)
    if not text_body:
        if specific_attachment_ids:
            raise ValueError("No extractable text found in the specified attachments.")
        else:
            raise ValueError("No extractable text found in this Material's attachments.")

    # ✅ IMPROVED: More specific prompt for unique titles
    system_prompt = (
        f"You are an AI tutor. Generate exactly {num_questions} multiple-choice questions "
        "based on the study material below. Create a SPECIFIC quiz title that reflects the "
        "actual content and topics covered (not generic titles like 'Quiz' or 'Test').\n\n"
        
        "**TITLE REQUIREMENTS:**\n"
        "- Include specific subject matter covered in the quiz\n"
        "- Add context like chapter/topic/difficulty if apparent\n"
        "- Use descriptive words that identify the content\n"
        "- Examples: 'Biology Quiz: Cellular Respiration', 'Chapter 5: World War II Events', 'Advanced Python: OOP Concepts'\n\n"
        
        "**RETURN ONLY RAW JSON - NO MARKDOWN, NO CODE BLOCKS, NO EXTRA TEXT.**\n"
        "Start your response immediately with { and end with }\n\n"
        "IMPORTANT: For 'correct_answer', provide the FULL TEXT of the correct choice, not just a letter.\n\n"
        "Expected JSON structure:\n"
        "{\n"
        "  \"title\": \"<specific descriptive quiz title>\",\n"
        "  \"description\": \"<description mentioning topics and difficulty>\",\n"
        "  \"questions\": [\n"
        "    {\n"
        "      \"question_text\": \"What is the main topic?\",\n"
        "      \"choices\": [\"First option\", \"Second option\", \"Third option\", \"Fourth option\"],\n"
        "      \"correct_answer\": \"Second option\"\n"
        "    },\n"
        f"    ... (exactly {num_questions} questions)\n"
        "  ]\n"
        "}\n\n"
        "REMEMBER: correct_answer must be the EXACT TEXT from one of the choices!\n\n"
        f"Material:\n\"\"\"\n{text_body}\n\"\"\"\n"
    )

    try:
        response = client.chat.completions.create(
            model="deepseek/deepseek-chat-v3-0324:free",
            messages=[{"role": "system", "content": system_prompt}],
        )
        raw = response.choices[0].message.content.strip()

        # ✅ Extract JSON from response (handles markdown code blocks)
        json_content = extract_json_from_response(raw)
        print(f"🔍 Quiz JSON content: {json_content[:200]}...")

        try:
            parsed = json.loads(json_content)
        except json.JSONDecodeError as e:
            print(f"❌ Quiz JSON parsing failed. Raw response: {raw}")
            raise ValueError(f"AI did not return valid JSON for quiz: {e}\n\nRaw response: {raw}")

        # ✅ Enhanced validation with better error messages
        title = parsed.get("title")
        description = parsed.get("description")
        questions = parsed.get("questions")

        print(f"🔍 Quiz validation - Title: {title}, Description: {description}, Questions count: {len(questions) if questions else 0}")

        if not isinstance(title, str) or not title.strip():
            raise ValueError("AI did not return a valid 'title' for the quiz.")
        if not isinstance(description, str):
            description = ""
        if not isinstance(questions, list) or len(questions) < 1:
            raise ValueError("AI JSON did not contain a valid 'questions' array.")

        validated = []
        for idx, qobj in enumerate(questions):
            try:
                # ✅ More robust validation
                if not isinstance(qobj, dict):
                    print(f"❌ Question #{idx+1} is not a dict: {qobj}")
                    raise ValueError(f"Question #{idx+1} is not a valid object: {qobj}")
                
                qt = qobj.get("question_text", "").strip()
                ch = qobj.get("choices")
                ca = qobj.get("correct_answer", "").strip()
                
                print(f"🔍 Question #{idx+1} - Text: '{qt[:50]}...', Choices: {len(ch) if ch else 0}, Correct: '{ca[:50]}...'")
                
                # Validate question text
                if not qt:
                    raise ValueError(f"Question #{idx+1} has empty question_text")
                
                # Validate choices
                if not isinstance(ch, list):
                    raise ValueError(f"Question #{idx+1} choices is not a list: {ch}")
                if len(ch) < 2:
                    raise ValueError(f"Question #{idx+1} needs at least 2 choices, got: {len(ch)}")
                
                # Clean up choices
                cleaned_choices = []
                for i, choice in enumerate(ch):
                    if not isinstance(choice, str):
                        raise ValueError(f"Question #{idx+1} choice #{i+1} is not a string: {choice}")
                    cleaned_choice = choice.strip()
                    if not cleaned_choice:
                        raise ValueError(f"Question #{idx+1} choice #{i+1} is empty")
                    cleaned_choices.append(cleaned_choice)
                
                # Validate correct answer
                if not ca:
                    raise ValueError(f"Question #{idx+1} has empty correct_answer")
                
                # ✅ Handle both letter-based and text-based correct answers
                if ca in ['A', 'B', 'C', 'D', 'E', 'F'] and len(cleaned_choices) >= ord(ca) - ord('A') + 1:
                    # Convert letter to actual choice text
                    choice_index = ord(ca.upper()) - ord('A')
                    ca = cleaned_choices[choice_index]
                    print(f"🔄 Converted letter '{qobj.get('correct_answer')}' to choice text: '{ca[:50]}...'")
                
                # Check if correct answer matches any choice
                if ca not in cleaned_choices:
                    print(f"❌ Question #{idx+1} correct answer '{ca}' not found in choices")
                    # Try to find a close match (case insensitive and trimmed)
                    ca_lower = ca.lower().strip()
                    matches = [choice for choice in cleaned_choices if choice.lower().strip() == ca_lower]
                    if matches:
                        ca = matches[0]  # Use the first match
                        print(f"✅ Found case-insensitive match: '{ca[:50]}...'")
                    else:
                        print(f"❌ Available choices:")
                        for i, choice in enumerate(cleaned_choices):
                            print(f"   {i+1}. {choice}")
                        raise ValueError(f"Question #{idx+1} correct_answer '{ca}' not found in choices")
                
                validated.append({
                    "question_text": qt,
                    "choices": cleaned_choices,
                    "correct_answer": ca,
                })
                print(f"✅ Question #{idx+1} validated successfully")
                
            except Exception as validation_error:
                print(f"❌ Question #{idx+1} validation failed: {validation_error}")
                raise ValueError(f"Question #{idx+1} validation failed: {validation_error}")

        print(f"✅ All {len(validated)} questions validated successfully")
        
        # ✅ Validate and improve title
        improved_title = validate_and_improve_title(title.strip(), 'quiz', material.title)
        
        return {
            "title": improved_title,
            "description": description.strip(),
            "questions": validated,
        }
    except Exception as e:
        print(f"❌ Quiz generation failed: {str(e)}")
        raise ValueError(f"Quiz generation failed: {str(e)}")

# ===== HELPER FUNCTIONS =====

def update_conversation_summary(conversation):
    """
    Update the conversation summary if needed.
    Call this from your view when processing messages.
    """
    if conversation.should_regenerate_summary():
        try:
            # Get messages to summarize (exclude the most recent ones)
            messages_to_summarize = conversation.messages[:-2] if len(conversation.messages) > 2 else conversation.messages
            
            # Generate new summary
            new_summary = generate_conversation_summary(
                messages_to_summarize, 
                conversation.summary_context
            )
            
            conversation.summary_context = new_summary
            conversation.reset_summary_counter()
            
            print(f"🔄 Generated new conversation summary: {new_summary[:100]}...")
            return True
            
        except Exception as e:
            print(f"⚠️ Failed to generate summary: {e}")
            return False
    
    return False